import os
import sys
import argparse
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Mm
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH

# ==========================================
# CONFIGURATION SECTION
# ==========================================
CONFIG = {
    "input_dir": "./input",
    "output_dir": "./output",
    "pptx_size": (Mm(190), Mm(270)),  # A4 Portrait (Width, Height)
    "file_prefix": "weekly_",
    "output_name_template": "monthly_report_{month}",
}

class MonthlyReportGenerator:
    def __init__(self, month_str):
        self.month_str = month_str # e.g., "2025-05"
        self.year, self.month = month_str.split('-')
        self.weekly_data = []
        
    def extract_text_from_pptx(self, file_path):
        """Extracts text from all slides in a pptx file."""
        prs = Presentation(file_path)
        extracted_text = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                extracted_text.append("\n".join(slide_text))
        return extracted_text

    def collect_weekly_reports(self):
        """Discovers and extracts data from all weekly reports in the input directory."""
        if not os.path.exists(CONFIG["input_dir"]):
            print(f"Error: Input directory '{CONFIG['input_dir']}' not found.")
            sys.exit(1)
            
        files = sorted([f for f in os.listdir(CONFIG["input_dir"]) if f.startswith(CONFIG["file_prefix"]) and f.endswith(".pptx")])
        
        if not files:
            print(f"Error: No files matching '{CONFIG['file_prefix']}*.pptx' found in '{CONFIG['input_dir']}'.")
            sys.exit(1)
            
        print(f"Found {len(files)} weekly reports. Extracting data...")
        for i, filename in enumerate(files):
            file_path = os.path.join(CONFIG["input_dir"], filename)
            content = self.extract_text_from_pptx(file_path)
            self.weekly_data.append({
                "week": i + 1,
                "filename": filename,
                "content": content
            })

    def generate_pptx(self):
        """Generates a portrait A4 PPTX report."""
        prs = Presentation()
        prs.slide_width, prs.slide_height = CONFIG["pptx_size"]
        
        # 1. Title Slide
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank slide for custom layout
        txBox = slide.shapes.add_textbox(Mm(10), Mm(100), Mm(170), Mm(50))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{self.year}년 {int(self.month)}월 월간보고"
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 2. TOC Slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Mm(10), Mm(10), Mm(170), Mm(20))
        title.text_frame.text = "목 차"
        title.text_frame.paragraphs[0].font.size = Pt(28)
        
        toc_items = ["1. 월간 핵심 요약"] + [f"{d['week']}주차 상세 내용" for d in self.weekly_data] + ["종합 정리"]
        for i, item in enumerate(toc_items):
            p = slide.shapes.add_textbox(Mm(20), Mm(40 + i*15), Mm(150), Mm(10)).text_frame.paragraphs[0]
            p.text = item
            p.font.size = Pt(18)

        # 3. Monthly Summary Slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Mm(10), Mm(10), Mm(170), Mm(20))
        title.text_frame.text = "월간 핵심 요약"
        
        body = slide.shapes.add_textbox(Mm(15), Mm(35), Mm(160), Mm(200))
        tf = body.text_frame
        tf.word_wrap = True
        for d in self.weekly_data:
            p = tf.add_paragraph()
            p.text = f"[{d['week']}주차 요약]"
            p.font.bold = True
            # Simple summary logic: first line of first content slide
            summary = d['content'][1].split('\n')[0] if len(d['content']) > 1 else "내용 없음"
            p = tf.add_paragraph()
            p.text = f"• {summary}"
            p.level = 1

        # 4. Weekly Detail Slides
        for d in self.weekly_data:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title = slide.shapes.add_textbox(Mm(10), Mm(10), Mm(170), Mm(20))
            title.text_frame.text = f"{d['week']}주차 상세 내용"
            
            body = slide.shapes.add_textbox(Mm(15), Mm(35), Mm(160), Mm(220))
            tf = body.text_frame
            tf.word_wrap = True
            for page in d['content']:
                for line in page.split('\n'):
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(12)

        # 5. Conclusion Slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Mm(10), Mm(10), Mm(170), Mm(20))
        title.text_frame.text = "종합 정리"
        
        body = slide.shapes.add_textbox(Mm(15), Mm(35), Mm(160), Mm(100))
        body.text_frame.text = "내용을 입력하세요."

        os.makedirs(CONFIG["output_dir"], exist_ok=True)
        out_path = os.path.join(CONFIG["output_dir"], CONFIG["output_name_template"].format(month=self.month_str.replace('-','')) + ".pptx")
        prs.save(out_path)
        print(f"PPTX saved to: {out_path}")

    def generate_docx(self):
        """Generates a DOCX report."""
        doc = Document()
        
        # 1. Cover
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(f"\n\n\n\n{self.year}년 {int(self.month)}월 월간보고")
        run.bold = True
        run.font.size = Pt(36)
        doc.add_page_break()

        # 2. TOC
        doc.add_heading("목 차", level=1)
        doc.add_paragraph("1. 월간 핵심 요약")
        for d in self.weekly_data:
            doc.add_paragraph(f"2.{d['week']} {d['week']}주차 상세 내용")
        doc.add_paragraph("3. 종합 의견")
        doc.add_page_break()

        # 3. Monthly Summary
        doc.add_heading("1. 월간 핵심 요약", level=1)
        for d in self.weekly_data:
            doc.add_heading(f"{d['week']}주차", level=2)
            summary = d['content'][1].split('\n')[0] if len(d['content']) > 1 else "내용 없음"
            doc.add_paragraph(summary, style='List Bullet')

        # 4. Weekly Details
        doc.add_heading("2. 주차별 상세 내용", level=1)
        for d in self.weekly_data:
            doc.add_heading(f"{d['week']}주차 상세 ({d['filename']})", level=2)
            for page in d['content']:
                doc.add_paragraph(page)
            doc.add_page_break()

        # 5. Conclusion
        doc.add_heading("3. 종합 의견", level=1)
        doc.add_paragraph("[내용 입력]")

        out_path = os.path.join(CONFIG["output_dir"], CONFIG["output_name_template"].format(month=self.month_str.replace('-','')) + ".docx")
        doc.save(out_path)
        print(f"DOCX saved to: {out_path}")

def main():
    parser = argparse.ArgumentParser(description="Generate Monthly Report from Weekly PPTXs")
    parser.add_argument("--month", required=True, help="Target month in YYYY-MM format")
    args = parser.parse_args()

    # Validate month format
    try:
        datetime.strptime(args.month, "%Y-%m")
    except ValueError:
        print("Error: Invalid month format. Please use YYYY-MM.")
        return

    generator = MonthlyReportGenerator(args.month)
    generator.collect_weekly_reports()

    print("\n생성할 파일 형식을 선택하세요:")
    print("1: PPTX 만 생성")
    print("2: DOCX 만 생성")
    print("3: 둘 다 생성")
    
    choice = input("선택 (1/2/3): ").strip()

    if choice in ['1', '3']:
        generator.generate_pptx()
    if choice in ['2', '3']:
        generator.generate_docx()
        
    if choice not in ['1', '2', '3']:
        print("잘못된 선택입니다. 종료합니다.")

if __name__ == "__main__":
    main()
