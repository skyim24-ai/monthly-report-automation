import os
import sys
import subprocess
import customtkinter as ctk
from tkinter import messagebox
import windnd
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt, Mm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import RGBColor as DocxRGBColor
from design_parser import DesignConfig

# ==========================================
# UPDATED LOGIC WITH DESIGN CONFIG
# ==========================================
class MonthlyReportGenerator:
    def __init__(self, year, month, design_cfg=None):
        self.year = year
        self.month = month
        self.month_str = f"{year}-{month:02d}"
        self.weekly_data = []
        self.cfg = design_cfg or DesignConfig()

    def extract_text_from_pptx(self, file_path):
        prs = Presentation(file_path)
        extracted_text = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                extracted_text.append("\n".join(slide_text))
        return extracted_text

    def process_files(self, file_paths):
        file_paths.sort()
        self.weekly_data = []
        for i, path in enumerate(file_paths):
            content = self.extract_text_from_pptx(path)
            self.weekly_data.append({
                "week": i + 1,
                "filename": os.path.basename(path),
                "content": content
            })

    def apply_pptx_font(self, run, size_key="body_size", is_bold=False, is_accent=False):
        style = self.cfg.get("pptx", "font_family")
        size = self.cfg.get("pptx", size_key)
        color = self.cfg.get("pptx", "accent_color" if is_accent else "title_color")
        
        run.font.name = style
        run.font.size = Pt(size)
        run.font.bold = is_bold
        run.font.color.rgb = RGBColor.from_string(color)

    def generate_pptx(self, output_dir):
        prs = Presentation()
        prs.slide_width, prs.slide_height = Mm(190), Mm(270)
        
        # Cover
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx = slide.shapes.add_textbox(Mm(10), Mm(100), Mm(170), Mm(50)).text_frame
        p = tx.paragraphs[0]
        p.text = f"{self.year}년 {self.month}월 월간보고"
        self.apply_pptx_font(p.runs[0], "title_size", True)
        p.alignment = PP_ALIGN.CENTER

        # Monthly Summary
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Mm(10), Mm(10), Mm(170), Mm(20))
        title_box.text_frame.text = "월간 핵심 요약"
        self.apply_pptx_font(title_box.text_frame.paragraphs[0].runs[0], "title_size", True)
        
        tf = slide.shapes.add_textbox(Mm(15), Mm(35), Mm(160), Mm(200)).text_frame
        tf.word_wrap = True
        for d in self.weekly_data:
            p = tf.add_paragraph()
            p.text = f"[{d['week']}주차 요약]"
            self.apply_pptx_font(p.runs[0], "body_size", True, True)
            
            summary = d['content'][1].split('\n')[0] if len(d['content']) > 1 else "내용 없음"
            p = tf.add_paragraph()
            p.text = f"• {summary}"
            self.apply_pptx_font(p.runs[0], "body_size")

        # Weekly Details
        for d in self.weekly_data:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Mm(10), Mm(10), Mm(170), Mm(20))
            title_box.text_frame.text = f"{d['week']}주차 상세"
            self.apply_pptx_font(title_box.text_frame.paragraphs[0].runs[0], "title_size", True)
            
            tf = slide.shapes.add_textbox(Mm(15), Mm(35), Mm(160), Mm(220)).text_frame
            tf.word_wrap = True
            for page in d['content']:
                for line in page.split('\n'):
                    p = tf.add_paragraph()
                    p.text = line
                    self.apply_pptx_font(p.runs[0], "body_size")

        out_name = f"monthly_report_{self.year}{self.month:02d}.pptx"
        path = os.path.join(output_dir, out_name)
        prs.save(path)
        return path

    def generate_docx(self, output_dir):
        doc = Document()
        
        # Helper to apply DOCX font
        def apply_docx_style(paragraph, size_key="body_size", is_bold=False):
            font_name = self.cfg.get("docx", "font_family")
            size = self.cfg.get("docx", size_key)
            color = self.cfg.get("docx", "text_color")
            
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(size)
                run.font.bold = is_bold
                run.font.color.rgb = DocxRGBColor.from_string(color)

        # Title
        p = doc.add_heading(f"{self.year}년 {self.month}월 월간보고", 0)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        apply_docx_style(p, "title_size", True)
        
        doc.add_heading("1. 월간 핵심 요약", level=1)
        for d in self.weekly_data:
            summary = d['content'][1].split('\n')[0] if len(d['content']) > 1 else "내용 없음"
            p = doc.add_paragraph(f"{d['week']}주차: {summary}", style='List Bullet')
            apply_docx_style(p, "body_size")

        doc.add_heading("2. 주차별 상세 내용", level=1)
        for d in self.weekly_data:
            p = doc.add_heading(f"{d['week']}주차 상세", level=2)
            apply_docx_style(p, "heading1_size", True)
            for page in d['content']:
                p = doc.add_paragraph(page)
                apply_docx_style(p, "body_size")
            doc.add_page_break()

        out_name = f"monthly_report_{self.year}{self.month:02d}.docx"
        path = os.path.join(output_dir, out_name)
        doc.save(path)
        return path

# ==========================================
# GUI WIDGET CLASS
# ==========================================
class MonthlyWidget(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title("월간 보고서 생성기 (Design.md 연동)")
        self.geometry("500x650")
        ctk.set_appearance_mode("dark")
        self.dropped_files = []
        self.design_cfg = DesignConfig()
        self.setup_ui()
        windnd.hook_dropfiles(self, self.on_drop)

    def setup_ui(self):
        self.label = ctk.CTkLabel(self, text="주간 보고서 취합 위젯", font=ctk.CTkFont(size=24, weight="bold"))
        self.label.pack(pady=20)

        # Design Status
        status_text = "design.md 적용 중" if os.path.exists("design.md") else "기본 디자인 사용 중"
        self.status_label = ctk.CTkLabel(self, text=status_text, font=ctk.CTkFont(size=12), text_color="gray70")
        self.status_label.pack(pady=(0, 10))

        self.drop_frame = ctk.CTkFrame(self, width=450, height=180, corner_radius=20, border_width=2, border_color="gray50")
        self.drop_frame.pack(pady=10, padx=25, fill="x")
        self.drop_frame.pack_propagate(False)

        self.drop_label = ctk.CTkLabel(self.drop_frame, text="PPTX 파일을 여기에 드래그하세요", font=ctk.CTkFont(size=16))
        self.drop_label.place(relx=0.5, rely=0.5, anchor="center")

        self.list_label = ctk.CTkLabel(self, text="선택된 파일 목록:", font=ctk.CTkFont(size=14, weight="bold"))
        self.list_label.pack(pady=(15, 5), padx=30, anchor="w")

        self.file_list_box = ctk.CTkTextbox(self, height=120, corner_radius=10)
        self.file_list_box.pack(pady=5, padx=25, fill="x")
        self.file_list_box.configure(state="disabled")

        self.btn_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.btn_frame.pack(pady=15, fill="x", padx=25)

        self.convert_btn = ctk.CTkButton(self.btn_frame, text="변환 시작", command=self.show_format_selection, state="disabled", height=45, font=ctk.CTkFont(size=16, weight="bold"))
        self.convert_btn.pack(side="left", expand=True, padx=5)

        self.clear_btn = ctk.CTkButton(self.btn_frame, text="초기화", command=self.clear_files, height=45, fg_color="gray40", hover_color="gray30")
        self.clear_btn.pack(side="right", expand=True, padx=5)

        self.edit_btn = ctk.CTkButton(self, text="디자인 수정", command=self.open_design_file, height=30, fg_color="transparent", border_width=1)
        self.edit_btn.pack(pady=5)

        self.progress = ctk.CTkProgressBar(self)
        self.progress.pack(pady=10, padx=25, fill="x")
        self.progress.set(0)

    def open_design_file(self):
        if os.path.exists("design.md"):
            os.startfile("design.md")
            messagebox.showinfo("알림", "디자인 파일을 열었습니다.\n수정 후 저장하시면 다음 변환 시 자동으로 반영됩니다.")
            # Reload immediately in case they already saved before clicking or we want to ensure latest
            self.design_cfg = DesignConfig() 
        else:
            messagebox.showerror("오류", "design.md 파일을 찾을 수 없습니다.")

    def reload_design(self):
        self.design_cfg = DesignConfig()
        status = "design.md 적용 중" if os.path.exists("design.md") else "기본 디자인 사용 중"
        self.status_label.configure(text=status)

    def on_drop(self, files):
        new_files = [f.decode('cp949') for f in files if f.decode('cp949').endswith('.pptx')]
        if not new_files:
            messagebox.showwarning("경고", "PPTX 파일만 드롭해주세요.")
            return
        self.dropped_files.extend(new_files)
        self.dropped_files = sorted(list(set(self.dropped_files)))
        self.update_file_list()
        self.convert_btn.configure(state="normal")
        self.drop_label.configure(text=f"{len(self.dropped_files)}개 파일 준비됨")

    def update_file_list(self):
        self.file_list_box.configure(state="normal")
        self.file_list_box.delete("1.0", "end")
        for i, f in enumerate(self.dropped_files):
            self.file_list_box.insert("end", f"{i+1}. {os.path.basename(f)}\n")
        self.file_list_box.configure(state="disabled")

    def clear_files(self):
        self.dropped_files = []
        self.update_file_list()
        self.convert_btn.configure(state="disabled")
        self.drop_label.configure(text="PPTX 파일을 여기에 드래그하세요")
        self.progress.set(0)

    def show_format_selection(self):
        dialog = ctk.CTkToplevel(self)
        dialog.title("출력 형식 선택")
        dialog.geometry("300x200")
        dialog.grab_set()
        label = ctk.CTkLabel(dialog, text="어떤 형식으로 저장할까요?", font=ctk.CTkFont(size=14))
        label.pack(pady=20)
        def select(choice):
            dialog.destroy()
            self.run_conversion(choice)
        ctk.CTkButton(dialog, text="PPTX", command=lambda: select(1)).pack(pady=5)
        ctk.CTkButton(dialog, text="DOCX", command=lambda: select(2)).pack(pady=5)
        ctk.CTkButton(dialog, text="둘 다 생성", command=lambda: select(3)).pack(pady=5)

    def run_conversion(self, choice):
        self.progress.set(0.3)
        self.update()
        now = datetime.now()
        gen = MonthlyReportGenerator(now.year, now.month, self.design_cfg)
        try:
            gen.process_files(self.dropped_files)
            output_dir = "./output"
            os.makedirs(output_dir, exist_ok=True)
            self.progress.set(0.6)
            self.update()
            if choice in [1, 3]: gen.generate_pptx(output_dir)
            if choice in [2, 3]: gen.generate_docx(output_dir)
            self.progress.set(1.0)
            messagebox.showinfo("완료", "월간 보고서 생성이 완료되었습니다.")
            subprocess.run(['explorer', os.path.realpath(output_dir)])
        except Exception as e:
            messagebox.showerror("오류", f"처리 중 오류가 발생했습니다: {e}")
            self.progress.set(0)

if __name__ == "__main__":
    app = MonthlyWidget()
    app.mainloop()
