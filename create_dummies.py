from pptx import Presentation
import os

def create_dummy_pptx(filename, week_num):
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"{week_num}주차 주간 회의자료"
    subtitle.text = f"작성일: 2025-05-{week_num*7:02d}"
    
    # Content Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "주요 성과"
    content.text = f"{week_num}주차 업무 실적 요약\n- 프로젝트 A 진행 완료\n- 이슈 B 해결"
    
    # Content Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "향후 계획"
    content.text = f"- 다음 주 {week_num+1}주차 업무 계획 수립\n- 고객사 미팅 예정"
    
    os.makedirs('input', exist_ok=True)
    prs.save(os.path.join('input', filename))
    print(f"Created {filename}")

if __name__ == "__main__":
    for i in range(1, 5): # Create 4 weeks by default
        create_dummy_pptx(f"weekly_{i:02d}.pptx", i)
    # Also create a 5th week to test variable count
    create_dummy_pptx("weekly_05.pptx", 5)
